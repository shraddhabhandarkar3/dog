# Import necessary libraries
import streamlit as st
from sql_module import (
    get_metadata_from_sql,
    update_metadata_steps,
    insert_evaluation,
    get_evaluations
)
from collections import Counter
from openai_module import send_to_openai
from aws_module import get_files_from_s3
import tempfile
import PyPDF2
import docx
import pandas as pd
import zipfile
import os
import openpyxl
import plotly.express as px
from dotenv import load_dotenv
import boto3
import easyocr
import base64
import re

# Load environment variables
load_dotenv()

# Define a consistent color palette
COLOR_PALETTE = {
    'green': '#2ca02c',
    'red': '#d62728',
    'blue': '#1f77b4',
    'orange': '#ff7f0e',
    'purple': '#9467bd',
    'cyan': '#17becf',
    'grey': '#7f7f7f'
}

# Supported file types
SUPPORTED_EXTENSIONS = [
    '.json', '.pdf', '.png', '.jpeg', '.jpg', '.txt',
    '.xlsx', '.csv', '.pptx', '.docx', '.py', '.zip', '.pdb'
]

# Initialize EasyOCR reader
reader = easyocr.Reader(['en'], gpu=False)

# Function to extract text from different file types
def extract_text_from_file(file_path):
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    try:
        if ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()

        elif ext == '.pdf':
            text = ""
            with open(file_path, 'rb') as f:
                reader_pdf = PyPDF2.PdfReader(f)
                for page in reader_pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            return text

        elif ext == '.docx':
            doc = docx.Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])

        elif ext == '.csv':
            df = pd.read_csv(file_path)
            return df.to_csv(index=False)

        elif ext == '.xlsx':
            try:
                df = pd.read_excel(file_path, engine='openpyxl')
                # Optionally display the dataframe
                # st.dataframe(df.head())

                # Convert the dataframe to a CSV string
                csv_representation = df.to_csv(index=False)
                return csv_representation

            except Exception as e:
                st.error(f"Error processing Excel file: {e}")
                return f"Error processing Excel file: {e}"

        elif ext in ['.png', '.jpg', '.jpeg']:
            try:
                # Use EasyOCR to extract text from the image
                result = reader.readtext(file_path, detail=0)
                text = ' '.join(result)
                return text
            except Exception as e:
                st.error(f"Error processing image file: {e}")
                return f"Error processing image file: {e}"

        elif ext == '.py':
            return extract_text_from_py(file_path)

        elif ext == '.zip':
            return extract_text_from_zip(file_path)

        elif ext == '.pdb':
            return extract_text_from_pdb(file_path)

        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)

        else:
            return f"Unsupported file type: {ext}"

    except Exception as e:
        return f"Error processing file: {e}"

# Helper function to extract text from .py files
def extract_text_from_py(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        with open(file_path, 'r', encoding='latin-1') as f:
            return f.read()
    except Exception as e:
        return f"Error processing .py file: {e}"

# Helper function to extract text from .zip files
def extract_text_from_zip(file_path):
    text = ""
    with zipfile.ZipFile(file_path, 'r') as zip_ref:
        with tempfile.TemporaryDirectory() as tmpdir:
            zip_ref.extractall(tmpdir)
            for root, _, files in os.walk(tmpdir):
                for file in files:
                    file_path_inner = os.path.join(root, file)
                    _, ext_inner = os.path.splitext(file_path_inner)
                    ext_inner = ext_inner.lower()
                    if ext_inner in SUPPORTED_EXTENSIONS:
                        extracted_text = extract_text_from_file(file_path_inner)
                        text += f"Extracted from {file}:\n{extracted_text}\n"
                    else:
                        text += f"Skipped unsupported file: {file}\n"
    return text

# Helper function to extract text from .pdb files
def extract_text_from_pdb(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        return f"Error processing .pdb file: {e}"

# Initialize Streamlit session state
def initialize_session_state():
    session_keys = [
        'comparison_result',
        'openai_response',
        'final_answer',
        'steps',
        'selected_task_id',
        'show_feedback_form',
        'show_rerun_button',
        'modified_steps',
        'awaiting_rerun_satisfaction',
        'awaiting_feedback'
    ]
    for key in session_keys:
        if key not in st.session_state:
            if key in ['comparison_result', 'show_feedback_form', 'show_rerun_button', 'awaiting_rerun_satisfaction', 'awaiting_feedback']:
                st.session_state[key] = False
            else:
                st.session_state[key] = ""

initialize_session_state()

# Custom CSS for styling
st.markdown("""
    <style>
    /* Global styles */
    body {
        font-family: 'Arial', sans-serif;
        font-size: 18px;
        color: #ffffff;
        background-color: #000000;
    }

    /* Header styles */
    h1 {
        font-size: 40px;
        font-weight: bold;
        color: #ffffff;
        text-align: center;
    }

    h2, h3 {
        color: #ffffff;
    }

    /* Paragraph and div styles */
    p, div {
        font-size: 18px;
        color: #ffffff;
    }

    /* Table styling */
    table {
        width: 100%;
        border-collapse: collapse;
    }

    th, td {
        padding: 14px;
        text-align: left;
        border-bottom: 1px solid #444444;
        font-size: 16px;
        color: #ffffff;
    }

    th {
        background-color: #333333;
        color: #ffffff;
    }

    /* Button styling */
    .stButton > button {
        background-color: #444444;
        color: #ffffff;
        border-radius: 10px;
        border: 1px solid #ffffff;
    }

    /* Selectbox styling */
    .stSelectbox > div {
        background-color: #333333;
        color: #ffffff;
        border: 1px solid #ffffff;
        border-radius: 5px;
    }

    /* Dataframe styling */
    .stDataFrame {
        background-color: #333333;
        color: #ffffff;
        border: 1px solid #ffffff;
        border-radius: 5px;
    }

    /* Sidebar styling */
    .sidebar .sidebar-content {
        background-color: #333333;
    }

    /* App container background */
    .stApp {
        background-color: #000000;
    }
    </style>
    """, unsafe_allow_html=True)

# Streamlit App Layout

# Sidebar for navigation and branding
with st.sidebar:
    st.image("https://streamlit.io/images/brand/streamlit-logo-primary-colormark-darktext.png", width=150)
    st.header("Navigation")
    st.write("Use the main page to process Task IDs and evaluate OpenAI responses.")
    st.markdown("---")
    st.header("Settings")
    # Future settings can be added here

# Main Page Header
st.markdown("""
    <h1>
        OpenAI Service Evaluation and Feedback Dashboard
    </h1>
    """, unsafe_allow_html=True)

st.markdown("---")

# Main Content Container
with st.container():
    # 1. Get metadata task_ids and questions from SQL Server
    metadata = get_metadata_from_sql()
    metadata_task_ids = [record['task_id'] for record in metadata]
    questions_dict = {record['task_id']: record['Question'] for record in metadata}
    final_answers_dict = {record['task_id']: record['Final answer'] for record in metadata}
    steps_dict = {record['task_id']: record['Steps'] for record in metadata}

    # 2. Get files from AWS S3 and create a mapping of task_ids to their files
    bucket_name = os.getenv('AWS_BUCKET')
    s3_files = get_files_from_s3(bucket_name)

    # Create a mapping from task_id to list of files (in case there are multiple files per task_id)
    task_files_mapping = {}
    for file_name in s3_files:
        # Extract task_id and file extension
        file_base_name, file_ext = os.path.splitext(file_name)
        if file_base_name not in task_files_mapping:
            task_files_mapping[file_base_name] = []
        task_files_mapping[file_base_name].append({'file_name': file_name, 'file_ext': file_ext.lower()})

    # 3. Include all task IDs from metadata
    all_task_ids = metadata_task_ids  # Include all task IDs, even those without files

    # 4. Select Task ID to query OpenAI
    selected_task_id = st.selectbox("Select a Task ID to Process", [""] + all_task_ids)
    st.session_state.selected_task_id = selected_task_id  # Store in session state

    if selected_task_id:
        # 5. Display the "Question" for the selected Task ID
        question = questions_dict[selected_task_id]
        st.markdown("### Question")
        st.write(question)

        # Initialize S3 client
        s3_client = boto3.client('s3')
        bucket_name = os.getenv('AWS_BUCKET')

        # Get the list of files associated with the selected task_id
        files_info = task_files_mapping.get(selected_task_id, [])

        # Initialize extracted_text
        extracted_text = ""

        # Process the files if any
        if files_info:
            for file_info in files_info:
                file_name = file_info['file_name']
                file_ext = file_info['file_ext']

                file_key = file_name  # The key in S3 is the file name

                # Retrieve and process the file from S3
                try:
                    if file_ext in SUPPORTED_EXTENSIONS:
                        # Download the file locally
                        with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tmp_file:
                            s3_client.download_fileobj(bucket_name, file_key, tmp_file)
                            tmp_file_path = tmp_file.name

                        # Extract text from the file
                        file_extracted_text = extract_text_from_file(tmp_file_path)
                        extracted_text += f"Extracted Text from {file_name}:\n{file_extracted_text}\n"

                    else:
                        st.write(f"Skipping unsupported file type for {file_name}")

                except Exception as e:
                    st.error(f"Error retrieving or processing the file {file_name} from S3: {e}")

        else:
            st.write("No associated files to process for this task.")

    # 6. Send question to OpenAI
    if st.button("Send Question to OpenAI", key="send_to_openai_button"):
        if selected_task_id:
            final_answer = final_answers_dict[selected_task_id]
            steps = steps_dict[selected_task_id]

            st.session_state.final_answer = final_answer
            st.session_state.steps = steps

            # Combine prompt with steps and extracted text if available
            prompt = f"Steps:\n{steps}\n\nQuestion:\n{question}\n"
            if extracted_text:
                prompt += f"\nExtracted Text:\n{extracted_text}"

            # Send to OpenAI and get response
            with st.spinner("Sending request to OpenAI..."):
                try:
                    result = send_to_openai(prompt)
                    st.session_state.openai_response = result

                    # Display OpenAI's response
                    st.markdown("### OpenAI's Response")
                    st.write(result)

                    # Display Final Answer
                    st.markdown("### Final Answer from Metadata")
                    st.write(final_answer)

                    # Reset rerun and feedback states
                    st.session_state.awaiting_rerun_satisfaction = False
                    st.session_state.awaiting_feedback = False
                except Exception as e:
                    st.error(f"An error occurred while communicating with OpenAI: {e}")

    # 7. User feedback on OpenAI response
    if (
        st.session_state.openai_response
        and not st.session_state.awaiting_feedback
        and not st.session_state.awaiting_rerun_satisfaction
    ):
        st.markdown("---")
        satisfaction = st.radio(
            "Is the OpenAI response satisfactory?", ("Yes", "No"), key="satisfaction_radio"
        )

        # Button for submitting feedback
        if st.button("Confirm Satisfaction", key="confirm_satisfaction_button"):
            if satisfaction == "No":
                # Trigger the feedback form display
                st.session_state.show_feedback_form = True
            elif satisfaction == "Yes":
                # Insert evaluation as correct without feedback
                st.markdown("**Submitting evaluation as correct...**")
                try:
                    insert_success = insert_evaluation(
                        task_id=selected_task_id,
                        is_correct=True,
                        user_feedback=None
                    )
                    if insert_success:
                        st.success("Evaluation recorded successfully.")
                        # Reset response and satisfaction state
                        st.session_state.openai_response = ""
                    else:
                        st.error("Failed to record evaluation.")
                except Exception as e:
                    st.error(f"An error occurred while recording evaluation: {e}")

    # 8. Show feedback form if the user selected "No"
    if st.session_state.show_feedback_form and not st.session_state.awaiting_feedback:
        st.markdown("---")
        st.markdown("### Modify Steps and Provide Feedback")
        # Display editable steps
        new_steps = st.text_area(
            "Edit the Steps below:",
            value=st.session_state.steps,
            height=200,
            key="edit_steps_textarea"
        )
        st.session_state.modified_steps = new_steps  # Store modified steps

        # Button to save modified steps and show rerun button
        if st.button("Save Modified Steps", key="save_modified_steps_button"):
            if st.session_state.modified_steps.strip() == "":
                st.error("Steps cannot be empty.")
            else:
                # Update the Steps in the metadata (SQL)
                try:
                    update_success = update_metadata_steps(
                        selected_task_id, st.session_state.modified_steps
                    )
                    if update_success:
                        st.success("Steps updated successfully.")
                        st.session_state.steps = st.session_state.modified_steps
                        st.session_state.show_feedback_form = False  # Hide feedback form after saving
                        st.session_state.show_rerun_button = True  # Show rerun model button
                    else:
                        st.error("Failed to update Steps. Please try again.")
                except Exception as e:
                    st.error(f"An error occurred: {e}")

    # 9. Show rerun model button if steps are modified
    if st.session_state.show_rerun_button and not st.session_state.awaiting_feedback:
        if st.button("Rerun Model", key="rerun_model_button"):
            # Combine prompt with modified steps and extracted text
            prompt = f"Steps:\n{st.session_state.modified_steps}\n\nQuestion:\n{question}\n"
            if extracted_text:
                prompt += f"\nExtracted Text:\n{extracted_text}"

            # Send to OpenAI and get new response
            with st.spinner("Rerunning the model with modified steps..."):
                try:
                    new_result = send_to_openai(prompt)
                    st.session_state.openai_response = new_result

                    # Display new OpenAI's response
                    st.markdown("### New OpenAI's Response")
                    st.write(new_result)

                    # Optionally, display Final Answer again
                    st.markdown("### Final Answer from Metadata")
                    st.write(final_answers_dict[selected_task_id])

                    # Reset rerun button and set flag to await satisfaction
                    st.session_state.show_rerun_button = False
                    st.session_state.awaiting_rerun_satisfaction = True
                except Exception as e:
                    st.error(f"An error occurred while rerunning the model: {e}")

    # 10. Handle Satisfaction Prompt After Rerun
    if (
        st.session_state.awaiting_rerun_satisfaction
        and not st.session_state.awaiting_feedback
    ):
        st.markdown("---")
        rerun_satisfaction = st.radio(
            "Is the new OpenAI response satisfactory?",
            ("Yes", "No"),
            key="rerun_satisfaction_radio"
        )

        if st.button("Confirm Rerun Satisfaction", key="confirm_rerun_satisfaction_button"):
            if rerun_satisfaction == "Yes":
                # Insert evaluation as correct without feedback
                st.markdown("**Submitting evaluation as correct...**")
                try:
                    insert_success = insert_evaluation(
                        task_id=selected_task_id,
                        is_correct=True,
                        user_feedback=None
                    )
                    if insert_success:
                        st.success("Evaluation recorded successfully.")
                        # Reset response and satisfaction state
                        st.session_state.openai_response = ""
                        st.session_state.awaiting_rerun_satisfaction = False
                    else:
                        st.error("Failed to record evaluation.")
                except Exception as e:
                    st.error(f"An error occurred while recording evaluation: {e}")
            elif rerun_satisfaction == "No":
                # Allow user to provide feedback
                st.session_state.awaiting_feedback = True
                st.session_state.awaiting_rerun_satisfaction = False

    # 11. Handle Feedback Submission After Rerun
    if st.session_state.awaiting_feedback:
        st.markdown("---")
        user_feedback = st.text_area(
            "Provide your feedback on the OpenAI response:",
            height=150,
            key="user_feedback_textarea"
        )

        if st.button("Submit Feedback", key="submit_feedback_button"):
            st.markdown("**Attempting to save feedback...**")
            if user_feedback.strip() == "":
                st.error("Feedback cannot be empty.")
            else:
                # Insert evaluation into Evaluations table with feedback
                try:
                    feedback_success = insert_evaluation(
                        task_id=selected_task_id,
                        is_correct=False,
                        user_feedback=user_feedback.strip()
                    )
                    if feedback_success:
                        st.success("Your feedback has been recorded.")
                        # Reset feedback and response states
                        st.session_state.awaiting_feedback = False
                        st.session_state.openai_response = ""
                    else:
                        st.error("Failed to record your feedback.")
                except Exception as e:
                    st.error(f"An error occurred: {e}")

    # 12. Generate Reports and Visualizations
    st.markdown("---")
    st.header("Evaluation Reports and Visualizations")

    # Fetch evaluations data
    evaluations = get_evaluations()
    if evaluations:
        eval_df = pd.DataFrame(evaluations)

        # Display basic metrics
        total_evaluations = len(eval_df)
        correct_answers = eval_df['is_correct'].sum()
        incorrect_answers = total_evaluations - correct_answers

        st.subheader("Summary Metrics")
        col1, col2, col3 = st.columns(3)
        col1.metric("Total Evaluations", total_evaluations, "")
        col2.metric("Correct Answers", correct_answers, f"{(correct_answers / total_evaluations) * 100:.2f}%")
        col3.metric("Incorrect Answers", incorrect_answers, f"{(incorrect_answers / total_evaluations) * 100:.2f}%")

        # Pie Chart for Correct vs Incorrect
        fig_pie = px.pie(
            names=['Correct', 'Incorrect'],
            values=[correct_answers, incorrect_answers],
            title='Distribution of OpenAI Responses',
            color=['Correct', 'Incorrect'],
            color_discrete_map={
                'Correct': COLOR_PALETTE['green'],
                'Incorrect': COLOR_PALETTE['red']
            }
        )
        fig_pie.update_layout(
            paper_bgcolor='#000000',  # Black background for the chart area
            plot_bgcolor='#000000',   # Black background for the plot area
            font=dict(color='#ffffff'),  # White font color for chart text
            title_font=dict(size=20, color='#ffffff')  # White for title
        )
        st.plotly_chart(fig_pie, use_container_width=True)

        # New Chart 1: Top 5 Most Common Feedback Themes
        # Simple text processing for feedback
        feedback_text = " ".join(eval_df['user_feedback'].dropna().tolist())
        # Remove non-alphabetic characters and lowercase
        words = re.findall(r'\b\w+\b', feedback_text.lower())
        # Define a list of stopwords (extend as needed)
        stopwords = set([
            'the', 'and', 'is', 'in', 'it', 'of', 'to', 'a', 'for', 'on',
            'with', 'as', 'this', 'that', 'but', 'be', 'have', 'are', 'was',
            'were', 'or', 'an', 'at', 'by', 'from', 'not', 'your', 'you'
        ])
        filtered_words = [word for word in words if word not in stopwords]
        word_counts = Counter(filtered_words)
        top_feedback = word_counts.most_common(5)
        feedback_labels, feedback_values = zip(*top_feedback) if top_feedback else ([], [])

        if feedback_labels:
            fig_feedback = px.bar(
                x=feedback_labels,
                y=feedback_values,
                title='Top 5 Most Common Feedback Themes',
                labels={'x': 'Feedback Theme', 'y': 'Frequency'},
                template='plotly_white',
                color=feedback_labels,
                color_discrete_map={
                    'incorrect': COLOR_PALETTE['red'],
                    'missing': COLOR_PALETTE['blue'],
                    'confusing': COLOR_PALETTE['orange'],
                    'too': COLOR_PALETTE['purple'],
                    'details': COLOR_PALETTE['cyan']
                }
            )
            fig_feedback.update_layout(
                paper_bgcolor='#000000',  # Black background for the chart area
                plot_bgcolor='#000000',   # Black background for the plot area
                font=dict(color='#ffffff'),  # White font color for chart text
                title_font=dict(size=20, color='#ffffff')  # White for title
            )
            st.plotly_chart(fig_feedback, use_container_width=True)
        else:
            st.write("No feedback available to display common themes.")

        # New Chart 2: Average Time Between Evaluations
        # Calculate time differences
        if 'evaluation_timestamp' in eval_df.columns:
            eval_df['evaluation_timestamp'] = pd.to_datetime(eval_df['evaluation_timestamp'])
            eval_df = eval_df.sort_values('evaluation_timestamp')
            eval_df['time_diff'] = eval_df['evaluation_timestamp'].diff().dt.total_seconds() / 60  # in minutes
            average_time_diff = eval_df['time_diff'].mean()
            st.subheader("Average Time Between Evaluations")
            if pd.notnull(average_time_diff):
                fig_time = px.histogram(
                eval_df,
                x='time_diff',
                nbins=30,
                title='Distribution of Time Differences Between Evaluations',
                labels={'time_diff': 'Time Difference (minutes)'},
                template='plotly_white',
                color_discrete_sequence=['#ff0000']  # Red color for bars
            )
                fig_time.update_layout(
                paper_bgcolor='#000000',  # Black background for the chart area
                plot_bgcolor='#000000',   # Black background for the plot area
                font=dict(color='#ffffff'),  # White font color for chart text
                title_font=dict(size=20, color='#ffffff')  # White for title
            )
                st.plotly_chart(fig_time, use_container_width=True)
            else:
                st.write("Not enough data to calculate time differences.")
        else:
            st.write("Evaluation timestamps not available for time-based analysis.")

        # Table of Evaluations
        st.subheader("Detailed Evaluations")
        st.dataframe(
            eval_df[['evaluation_id', 'task_id', 'is_correct', 'user_feedback', 'evaluation_timestamp']],
            height=300
        )
    else:
        st.write("No evaluations recorded yet.")
